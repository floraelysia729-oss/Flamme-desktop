"""
LLM Wiki 摄入管道 — PPT/PDF → Markdown

用法:
  python scripts/ingest.py <file_or_dir> [--name 名称] [--ppt2pdf]

示例:
  python scripts/ingest.py "E:/课件/矩阵论.pptx" --name 矩阵论
  python scripts/ingest.py "D:/notebook/数字系统设计"
  python scripts/ingest.py chap1.pdf --name "计算机系统"
"""

import argparse
import os
import sys
import tempfile
from pathlib import Path

# 确保 UTF-8 输出 + 禁用缓冲（管道/后台时也能看到进度）
sys.stdout.reconfigure(encoding="utf-8", line_buffering=True)

from src.scripts import VAULT
from src.tools.paths import converted_dir


# ── PPT → PDF (Windows, via PowerPoint COM) ──────────────────────────

def pptx_to_pdf(pptx_path: Path, out_dir: Path) -> Path:
    """用 PowerPoint COM 接口将 .pptx 转为 .pdf"""
    import comtypes.client

    pdf_path = out_dir / (pptx_path.stem + ".pdf")
    if pdf_path.exists():
        print(f"  [skip] PDF already exists: {pdf_path.name}")
        return pdf_path

    print(f"  [ppt→pdf] {pptx_path.name} ...")
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1
    abs_pptx = str(pptx_path.resolve())
    abs_pdf = str(pdf_path.resolve())
    deck = powerpoint.Presentations.Open(abs_pptx)
    deck.SaveAs(abs_pdf, 32)  # 32 = ppSaveAsPDF
    deck.Close()
    powerpoint.Quit()
    print(f"  [done] → {pdf_path.name}")
    return pdf_path


# ── PPT → Markdown (python-pptx, 直接提取) ──────────────────────────

def _table_to_md(table) -> str:
    """将 python-pptx Table 转为 Markdown 表格"""
    rows = []
    for row in table.rows:
        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
        rows.append("| " + " | ".join(cells) + " |")
    if len(rows) >= 1:
        header = rows[0]
        sep = "|" + "|".join("---" for _ in table.columns) + "|"
        return header + "\n" + sep + "\n" + "\n".join(rows[1:])
    return ""


def _format_text_frame(text_frame) -> str:
    """提取 text_frame 内容，保留粗体和列表层级"""
    lines = []
    for para in text_frame.paragraphs:
        if not para.text.strip():
            continue
        level = para.level
        prefix = "  " * level + "- " if level > 0 else "- "
        runs_text = []
        for run in para.runs:
            if run.font.bold:
                runs_text.append(f"**{run.text}**")
            else:
                runs_text.append(run.text)
        line = "".join(runs_text).strip()
        if line:
            if line.startswith(('- ', '• ', '* ')):
                indent = "  " * level
                lines.append(f"{indent}{line}")
            else:
                lines.append(f"{prefix}{line}")
    return "\n".join(lines)


def pptx_to_markdown(pptx_path: Path) -> str:
    """用 python-pptx 直接提取 PPT 文字和表格"""
    from pptx import Presentation

    print(f"  [extract] {pptx_path.name} ...")
    prs = Presentation(str(pptx_path))
    lines = [f"# {pptx_path.stem}\n"]

    for i, slide in enumerate(prs.slides, 1):
        slide_title = ""
        body_parts = []
        notes_text = ""

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()
                if not text:
                    continue
                if not slide_title:
                    slide_title = text
                else:
                    body_parts.append(_format_text_frame(shape.text_frame))

            elif shape.has_table:
                md_table = _table_to_md(shape.table)
                if md_table:
                    body_parts.append(md_table)

        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            notes_text = notes_frame.text.strip()

        if slide_title or body_parts:
            lines.append(f"\n## Slide {i}: {slide_title}\n")
            for part in body_parts:
                lines.append(part)
                lines.append("")
            if notes_text:
                lines.append(f"> **备注**: {notes_text}")

    return "\n".join(lines)


# ── PDF → Markdown (MinerU) ─────────────────────────────────────────

def pdf_to_markdown(pdf_path: Path) -> str:
    """用 MinerU 云端 API 解析 PDF 为 Markdown"""
    from src.tools.pdf_parser import PDFParserTool
    token = os.environ.get("MINERU_API_TOKEN", "")
    if not token:
        raise RuntimeError("MINERU_API_TOKEN 未配置")
    tool = PDFParserTool(api_token=token)
    result = tool.execute({"path": str(pdf_path)})
    if result.is_error:
        raise RuntimeError(f"MinerU 解析失败: {result.error}")
    return result.data["markdown"]


# ── Frontmatter ────────────────────────────────────────────────────

def make_frontmatter(title: str, source: str, level: str, tags: list[str] = None) -> str:
    from datetime import date

    today = date.today().isoformat()
    tag_str = str(tags or []).replace("'", '"')
    return f"""---
title: "{title}"
date: {today}
source: "{source}"
level: {level}
tags: {tag_str}
related: []
status: stable
---

"""


# ── 主流程 ─────────────────────────────────────────────────────────

def process_file(filepath: Path, name: str = None, ppt2pdf_flag: bool = False):
    """处理单个文件 — 输出到 {课程名}/.flamme/converted/"""
    ext = filepath.suffix.lower()
    dir_name = name or filepath.stem
    source_dir = VAULT / dir_name
    source_dir.mkdir(parents=True, exist_ok=True)

    file_stem = filepath.stem

    if ext == ".pdf":
        md_text = pdf_to_markdown(filepath)
        source = str(filepath.relative_to(VAULT)) if filepath.is_relative_to(VAULT) else filepath.name

        md_text = make_frontmatter(file_stem, source, "source") + md_text
        out_md = converted_dir(source_dir) / f"{file_stem}.md"
        out_md.write_text(md_text, encoding="utf-8")
        print(f"  → {out_md}")
        return out_md

    elif ext == ".pptx":
        # 方式1: 直接用 python-pptx 提取文字
        md_text = pptx_to_markdown(filepath)
        source = filepath.name

        # 方式2: 如果指定 --ppt2pdf，额外通过 MinerU 解析 PDF 版本
        if ppt2pdf_flag:
            try:
                pdf_path = pptx_to_pdf(filepath, converted_dir(source_dir))
                pdf_text = pdf_to_markdown(pdf_path)
                md_text = f"{md_text}\n\n---\n\n# PDF 提取补充内容\n\n{pdf_text}"
                source = f"{filepath.name} (pptx+pdf)"
            except Exception as e:
                print(f"  [warn] ppt2pdf failed: {e}")

        md_text = make_frontmatter(file_stem, source, "source") + md_text
        out_md = converted_dir(source_dir) / f"{file_stem}.md"
        out_md.write_text(md_text, encoding="utf-8")
        print(f"  → {out_md}")
        return out_md

    elif ext == ".md":
        import shutil
        out_md = converted_dir(source_dir) / filepath.name
        shutil.copy2(filepath, out_md)
        print(f"  → copied to {out_md}")
        return out_md

    else:
        print(f"  [skip] unsupported format: {ext}")
        return None


def main():
    parser = argparse.ArgumentParser(description="LLM Wiki 摄入管道")
    parser.add_argument("path", help="文件或目录路径")
    parser.add_argument("--name", help="自定义目录名（默认用文件名）")
    parser.add_argument("--ppt2pdf", action="store_true", help="PPT 额外转 PDF 用 MinerU 解析（需安装 PowerPoint）")
    args = parser.parse_args()

    path = Path(args.path)
    if not path.exists():
        print(f"Error: {path} not found")
        sys.exit(1)

    # 收集要处理的文件
    supported = {".pdf", ".pptx"}
    if path.is_dir():
        files = sorted(p for p in path.rglob("*")
                       if p.suffix.lower() in supported and not p.name.startswith("~$"))
    else:
        files = [path]

    if not files:
        print("No supported files found")
        sys.exit(1)

    # 检查 MinerU token
    if not os.environ.get("MINERU_API_TOKEN"):
        print("Warning: MINERU_API_TOKEN not set. PDF parsing will fail.")
        print("  Set it via: export MINERU_API_TOKEN=your-token")
        print()

    print(f"Found {len(files)} file(s) to process\n")

    results = []
    for f in files:
        print(f"[*] {f.name}")
        result = process_file(f, args.name, args.ppt2pdf)
        if result:
            results.append(result)
        print()

    print(f"Done. {len(results)} file(s) processed.")
    if results:
        print("Output:")
        for r in results:
            print(f"  {r}")


if __name__ == "__main__":
    main()
